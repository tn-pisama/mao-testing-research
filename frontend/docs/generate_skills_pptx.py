#!/usr/bin/env python3
"""Generate editable PowerPoint presentation for Claude Skills strategy."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Colors
DARK_BG = (15, 23, 42)
PURPLE = (167, 139, 250)
LIGHT_PURPLE = (196, 181, 253)
WHITE = (241, 245, 249)

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_slide(prs, title, subtitle=None, body=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(*DARK_BG)
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.8), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = rgb(*PURPLE)
    p.font.name = "Arial"

    y = 1.3

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(y), Inches(11.8), Inches(0.5))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = rgb(*LIGHT_PURPLE)
        p.font.name = "Arial"
        y = 2.0

    if body:
        body_box = slide.shapes.add_textbox(Inches(0.75), Inches(y), Inches(11.8), Inches(5))
        body_box.text_frame.word_wrap = True
        first = True
        for line in body.split("\n"):
            if first:
                p = body_box.text_frame.paragraphs[0]
                first = False
            else:
                p = body_box.text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = rgb(*WHITE)
            p.font.name = "Arial"
            p.space_after = Pt(6)

    return slide

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = [
        ("PISAMA x Claude Skills", "Agent Forensics Powered by Anthropic's Extensibility Platform",
         "pisama.ai"),

        ("The Opportunity", "Three Strategic Goals",
         """1. BUILD - Use Claude Skills to accelerate PISAMA development

2. DISTRIBUTE - Ship PISAMA capabilities as Claude Skills

3. COLLECT - Access user skill execution traces for analysis


One platform, three integration layers"""),

        ("What Are Claude Skills?", "Claude's Native Extensibility Mechanism",
         """• Markdown instruction files that extend Claude Code capabilities
• Executed within Claude's environment with full tool access
• Users can install, share, and customize skills
• MCP (Model Context Protocol) servers provide integrations

KEY INSIGHT:
Skills run inside Claude = perfect observation point for agent behavior"""),

        ("Goal 1: BUILD with Skills", "Internal Development Acceleration",
         """pisama-architect
    Design detection algorithms, suggest patterns

detection-designer
    Create new failure detectors from examples

sdk-generator
    Auto-generate SDK code for new platforms

test-synthesizer
    Generate test cases from failure patterns

BENEFIT: 3-5x faster feature development"""),

        ("Goal 2: DISTRIBUTE as Skills", "PISAMA Capabilities as Installable Skills",
         """/install pisama-diagnose

PRODUCT SKILLS PORTFOLIO:
• pisama-learn - Interactive tutorials on agent failure patterns
• pisama-diagnose - Analyze current project for potential issues
• pisama-fix - Apply recommended fixes automatically
• pisama-review - Pre-commit review for agent anti-patterns
• pisama-monitor - Set up observability in existing projects"""),

        ("Distribution Advantage", "Why Skills as Distribution Channel?",
         """TRADITIONAL SAAS          SKILLS DISTRIBUTION
User visits website       User types /install
Sign up flow              Already authenticated
Learn new UI              Uses familiar Claude
Integration work          Works in environment
Monthly subscription      Usage-based


RESULT: Zero-friction adoption in developer workflow"""),

        ("Goal 3: COLLECT Traces", "The Core Innovation: MCP Trace Observer",
         """┌─────────────────────────────────────────┐
│           Claude Code Session           │
│  ┌─────────────────────────────────┐   │
│  │     User's Skill Execution      │   │
│  └─────────────┬───────────────────┘   │
│       ┌────────▼────────┐              │
│       │  PISAMA MCP     │              │
│       │  Trace Observer │              │
│       └────────┬────────┘              │
└────────────────┼────────────────────────┘
        ┌────────▼────────┐
        │  PISAMA Cloud   │
        └─────────────────┘"""),

        ("What We Capture", "Trace Data Points",
         """• Tool Calls - File reads, writes, bash commands, searches
• Decision Points - When Claude chooses between options
• Iteration Patterns - Retry loops, error corrections
• Context Switches - Topic changes, goal modifications
• Timing - Duration of each phase

NOT CAPTURED:
File contents, secrets, personal data (privacy-first design)"""),

        ("Privacy Framework", "User Trust is Non-Negotiable",
         """Opt-in Only
    Explicit consent before collection

Minimal Data
    Hashed identifiers, no raw content

User Ownership
    Export, delete, view your data

Transparency
    Open-source collection logic

Local-First
    Analysis can run on-device

CERTIFICATION: SOC2, GDPR compliant from day one"""),

        ("Detection Capabilities", "What PISAMA Identifies From Traces",
         """• Infinite Loops - Repeated tool call patterns
• State Corruption - Semantic drift in context
• Goal Abandonment - Incomplete task chains
• Resource Exhaustion - Token/time budget overruns
• Coordination Failures - Multi-agent deadlocks

ACCURACY TARGET: 95%+ precision on known failure modes"""),

        ("User Value Proposition", "For Claude Code Users",
         """BEFORE PISAMA:
"Why did my 2-hour Claude session fail to complete the task?"

AFTER PISAMA:
"PISAMA detected a loop at minute 47 where Claude kept
re-reading the same 3 files. Suggested fix: add explicit
progress checkpoints. One-click apply."

TIME SAVED: 30-60 minutes per failed session"""),

        ("Implementation Phases", "16-Week Roadmap",
         """PHASE 1: Foundation (Weeks 1-4)
    MCP server, basic trace capture

PHASE 2: Detection (Weeks 5-8)
    Loop & corruption detectors

PHASE 3: Skills v1 (Weeks 9-12)
    pisama-diagnose, pisama-fix

PHASE 4: Scale (Weeks 13-16)
    Cloud platform, dashboards"""),

        ("Technical Architecture", None,
         """┌──────────────────────────────────────────────────────┐
│                    User Environment                   │
│  ┌────────────┐  ┌────────────┐  ┌────────────────┐ │
│  │ Claude Code│  │ PISAMA MCP │  │ PISAMA Skills  │ │
│  │  Session   │◄─┤  Observer  │  │ (diagnose/fix) │ │
│  └────────────┘  └─────┬──────┘  └───────┬────────┘ │
└────────────────────────┼─────────────────┼──────────┘
              ┌──────────▼─────────────────▼──────────┐
              │           PISAMA Cloud                 │
              │  ┌─────────┐ ┌──────────┐ ┌────────┐  │
              │  │ Traces  │ │ Analysis │ │ Fixes  │  │
              │  └─────────┘ └──────────┘ └────────┘  │
              └────────────────────────────────────────┘"""),

        ("Business Model", "Revenue Streams",
         """Free ($0)
    Basic diagnostics, community patterns

Pro ($29/mo)
    Full analysis, custom detectors

Team ($99/mo)
    Team dashboards, shared patterns, API

Enterprise (Custom)
    On-prem, SLA, dedicated support

TARGET: $1M ARR in 18 months"""),

        ("Competitive Moat", "Why PISAMA Wins",
         """1. First Mover - No dedicated Claude Skills forensics tool exists

2. Data Network Effect - More users = better pattern detection

3. Platform Integration - Native to Claude, not bolted on

4. Open Core - Community contributes detection patterns

5. AI-Native - Built by AI developers, for AI developers"""),

        ("Success Metrics", "KPIs to Track",
         """                        6-MONTH     12-MONTH
Skills Installs          5,000       50,000
Monthly Active Users     1,000       10,000
Traces Analyzed          100K        2M
Detection Accuracy       90%         95%
Paid Conversion          3%          5%"""),

        ("Why Now?", "Market Timing",
         """• Claude Code adoption accelerating rapidly
• No mature observability for AI coding agents
• Anthropic actively promoting MCP ecosystem
• Developer pain points are acute and unaddressed
• Skills marketplace being built out

WINDOW: 12-18 months before major players enter"""),

        ("Next Steps", "Immediate Actions",
         """1. Complete MCP server prototype (Week 1-2)
2. Build trace capture pipeline (Week 2-3)
3. Develop first detection algorithms (Week 3-4)
4. Create pisama-diagnose skill (Week 5-6)
5. Private beta with 50 users (Week 7-8)
6. Public launch on Skills marketplace (Week 12)"""),

        ("The Ask", "What We Need",
         """DEVELOPMENT
    2 engineers, 16 weeks

INFRASTRUCTURE
    Cloud hosting, trace storage

PARTNERSHIPS
    Early access to Anthropic Skills APIs

VALIDATION
    50 beta testers from Claude Code power users

INVESTMENT: Seeking seed round to accelerate"""),

        ("PISAMA: Agent Forensics", None,
         """"Find out why your AI agent failed. Fix it automatically."

• Build faster with internal skills
• Distribute through Claude's native platform
• Collect traces to power intelligent diagnostics

Website: pisama.ai"""),
    ]

    for args in slides:
        add_slide(prs, *args)

    return prs

if __name__ == "__main__":
    prs = create_presentation()
    prs.save("PISAMA_Skills_Editable.pptx")
    print("Created: PISAMA_Skills_Editable.pptx")
