#!/usr/bin/env python3
"""Generate editable PowerPoint presentation for partner."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Colors as hex for RGB
DARK_BG = (15, 23, 42)  # #0f172a
PURPLE = (167, 139, 250)  # #a78bfa
LIGHT_PURPLE = (196, 181, 253)  # #c4b5fd
WHITE = (241, 245, 249)  # #f1f5f9

def rgb_color(r, g, b):
    """Helper to create RGB color."""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_content = [
        # Slide 1: Title
        {
            "title": "PISAMA",
            "subtitle": "What I've Been Building",
            "body": "A tool that helps AI systems work better",
            "type": "title"
        },
        # Slide 2: What are AI Agents
        {
            "title": "First, What Are AI Agents?",
            "body": """You know ChatGPT - you ask it questions and it answers.

AI Agents are the next step: AI that can actually DO things.

• Book your flights
• Write and send emails
• Research and summarize documents
• Build websites
• Manage your calendar

They're like robot assistants that can take action, not just chat.""",
            "type": "content"
        },
        # Slide 3: The Problem
        {
            "title": "The Big Problem",
            "subtitle": "AI Agents Break. A Lot.",
            "body": """Imagine asking your assistant to book a trip, and they:

• Keep checking the same flight over and over (infinite loop)
• Forget they're your assistant and start acting weird
• Get stuck waiting for approval that never comes
• Book the wrong dates because they misread your request

Nobody knows why these failures happen.

Developers spend HOURS reading logs trying to figure it out.""",
            "type": "content"
        },
        # Slide 4: What I Built
        {
            "title": "What I Built",
            "subtitle": "A Doctor for AI Agents",
            "body": """PISAMA watches AI agents work and:

1. DETECTS when something goes wrong

2. EXPLAINS exactly why it failed

3. SUGGESTS how to fix it automatically

Think of it like a health monitoring system for robot assistants.""",
            "type": "content"
        },
        # Slide 5: How It Works
        {
            "title": "How It Works (Simply)",
            "body": """Your AI Agent (doing tasks)
        ↓
    PISAMA Monitor  ←  Watches everything
        ↓
    Problem Found!  ←  "Loop detected!"
        ↓
    Here's The Fix  ←  "Add a 3-try limit"


Simple: Watch → Detect → Explain → Fix""",
            "type": "content"
        },
        # Slide 6: What Can It Catch
        {
            "title": "What Can It Catch?",
            "subtitle": "26 Different Types of Problems",
            "body": """• Infinite Loops - The AI keeps repeating the same action

• Memory Problems - The AI forgets important information

• Identity Crisis - The AI forgets what role it's supposed to play

• Stuck Waiting - Multiple AIs waiting on each other forever

• Making Things Up - The AI invents facts that aren't true

• Security Issues - Someone trying to trick the AI""",
            "type": "content"
        },
        # Slide 7: Why Does This Matter
        {
            "title": "Why Does This Matter?",
            "subtitle": "AI Agents Are Becoming Huge",
            "body": """Companies are building AI agents to:
• Handle customer support
• Process legal documents
• Manage supply chains
• Write software
• Run entire businesses

The problem: When these agents fail, businesses lose money.

The opportunity: Nobody has built good tools to fix this. Until now.""",
            "type": "content"
        },
        # Slide 8: The Market
        {
            "title": "The Market",
            "subtitle": "This Is a Real Business Opportunity",
            "body": """• $50+ billion spent on AI development in 2024

• Companies building AI agents have no good debugging tools

• Existing tools weren't designed for AI agents

• First-mover advantage - we're early

Think about it like the early days of the internet:
everyone was building websites, but there were no good tools yet.""",
            "type": "content"
        },
        # Slide 9: What I Actually Built
        {
            "title": "What I Actually Built",
            "subtitle": "It's Real and Working",
            "body": """• Website Dashboard - See all your AI agents in one place

• Detection Engine - Finds 26 types of problems automatically

• Fix Suggester - Tells you exactly how to repair issues

• Connectors - Works with 4 popular AI frameworks

• Command Line Tool - For developers who prefer typing

NOT a prototype - it's production-ready software.""",
            "type": "content"
        },
        # Slide 10: The Dashboard
        {
            "title": "The Dashboard",
            "subtitle": "Where You See Everything",
            "body": """Imagine a control room where you can:

• See all your AI agents running
• Watch for problems in real-time
• Get alerts when something breaks
• Click to see exactly what went wrong
• Get suggested fixes instantly

It's like a mission control for AI systems.""",
            "type": "content"
        },
        # Slide 11: Demo
        {
            "title": "A Quick Demo",
            "subtitle": "You Can Try It Right Now",
            "body": """pisama.ai/demo

1. Pick a failure scenario (like "Infinite Loop")

2. Watch the AI agents try to work

3. See PISAMA catch the problem

4. View the explanation and fix

No account needed - just click and watch.""",
            "type": "content"
        },
        # Slide 12: How Long
        {
            "title": "How Long Did This Take?",
            "subtitle": "Months of Work",
            "body": """This involved building:

• Complex algorithms to detect subtle failures
• A database system to store all the information
• A beautiful web interface
• Tools that connect to other AI systems
• Security and user accounts
• Documentation for developers

This is not a weekend project - it's a real product.""",
            "type": "content"
        },
        # Slide 13: Technical (Brief)
        {
            "title": "The Technical Stuff (Briefly)",
            "subtitle": "You Don't Need to Understand This",
            "body": """But just so you know what's under the hood:

• Backend: Python + FastAPI (modern, fast)
• Database: PostgreSQL (reliable, used by Netflix)
• Frontend: Next.js + React (what Facebook uses)
• AI Detection: Custom algorithms + embeddings

It's built with professional-grade technology.""",
            "type": "content"
        },
        # Slide 14: Competition
        {
            "title": "Competition",
            "subtitle": "Who Else Is Doing This?",
            "body": """LangSmith - General AI logging
    → Not focused on agent failures

Arize - ML monitoring
    → Not for agent failures

Braintrust - AI testing
    → No self-healing

Nobody is specifically solving the "AI agent failure" problem.

That's our unique position.""",
            "type": "content"
        },
        # Slide 15: Business Model
        {
            "title": "Business Model",
            "subtitle": "How It Makes Money",
            "body": """Free ($0) - Hobbyists, trying it out

Pro ($29/month) - Individual developers

Team ($99/month) - Small companies

Enterprise (Custom) - Large corporations


Standard software-as-a-service model.
Recurring revenue that grows with usage.""",
            "type": "content"
        },
        # Slide 16: What's Next
        {
            "title": "What's Next",
            "subtitle": "The Roadmap",
            "body": """Coming Soon:
• More AI framework connections
• Automatic fix application (one-click repair)
• Team collaboration features

Future:
• Partnerships with AI companies
• Enterprise sales
• Potentially: acquisition target for bigger players""",
            "type": "content"
        },
        # Slide 17: Why I'm Excited
        {
            "title": "Why I'm Excited",
            "subtitle": "This Solves a Real Problem",
            "body": """Every developer building AI agents today struggles with debugging.

I've felt this pain myself.
Hours lost trying to understand why an AI failed.

PISAMA makes that 10x faster.

When AI agents become as common as websites,
everyone will need this.""",
            "type": "content"
        },
        # Slide 18: Vision
        {
            "title": "The Vision",
            "subtitle": "Where This Could Go",
            "body": """Short term:
Tool for developers building AI agents

Medium term:
Standard part of every AI development workflow

Long term:
The essential infrastructure for the AI agent economy

Like how every website needs hosting,
every AI agent will need monitoring.""",
            "type": "content"
        },
        # Slide 19: Thank You
        {
            "title": "Thank You",
            "subtitle": "For Your Support",
            "body": """Building this has meant late nights and weekends.

Your patience and encouragement made it possible.

This is just the beginning.""",
            "type": "thank_you"
        },
        # Slide 20: Questions
        {
            "title": "Questions?",
            "subtitle": "Happy to Explain Anything",
            "body": """Website: pisama.ai

Live Demo: pisama.ai/demo

Or just ask me!""",
            "type": "content"
        },
    ]

    for slide_data in slides_content:
        create_slide(prs, slide_data)

    return prs


def create_slide(prs, data):
    """Create a slide with the given content."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = rgb_color(*DARK_BG)
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.5), Inches(11.8), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = data["title"]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = rgb_color(*PURPLE)
    title_para.font.name = "Arial"

    y_offset = 1.5

    # Subtitle if present
    if "subtitle" in data:
        sub_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.4), Inches(11.8), Inches(0.6)
        )
        sub_frame = sub_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = data["subtitle"]
        sub_para.font.size = Pt(28)
        sub_para.font.color.rgb = rgb_color(*LIGHT_PURPLE)
        sub_para.font.name = "Arial"
        y_offset = 2.2

    # Body content
    if "body" in data:
        body_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(y_offset), Inches(11.8), Inches(5)
        )
        body_frame = body_box.text_frame
        body_frame.word_wrap = True

        # Split by lines and add paragraphs
        lines = data["body"].split("\n")
        first = True
        for line in lines:
            if first:
                para = body_frame.paragraphs[0]
                first = False
            else:
                para = body_frame.add_paragraph()

            para.text = line
            para.font.size = Pt(22)
            para.font.color.rgb = rgb_color(*WHITE)
            para.font.name = "Arial"
            para.space_after = Pt(8)

    return slide


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "PISAMA_Partner_Editable.pptx"
    prs.save(output_path)
    print(f"Created: {output_path}")
