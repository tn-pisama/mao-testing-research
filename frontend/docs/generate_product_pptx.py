#!/usr/bin/env python3
"""Generate editable PowerPoint presentation for PISAMA product."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Colors
DARK_BG = (15, 23, 42)
PURPLE = (167, 139, 250)
LIGHT_PURPLE = (196, 181, 253)
WHITE = (241, 245, 249)
CYAN = (34, 211, 238)

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_slide(prs, title, subtitle=None, body=None, table_data=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(*DARK_BG)
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.8), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = rgb(*PURPLE)
    p.font.name = "Arial"

    y = 1.3

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(y), Inches(11.8), Inches(0.5))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = rgb(*LIGHT_PURPLE)
        p.font.name = "Arial"
        y = 2.0

    # Body
    if body:
        body_box = slide.shapes.add_textbox(Inches(0.75), Inches(y), Inches(11.8), Inches(5))
        body_box.text_frame.word_wrap = True
        first = True
        for line in body.split("\n"):
            if first:
                p = body_box.text_frame.paragraphs[0]
                first = False
            else:
                p = body_box.text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = rgb(*WHITE)
            p.font.name = "Arial"
            p.space_after = Pt(6)

    return slide

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = [
        ("PISAMA", "Agent Forensics Platform", "Find out why your AI agent failed. Fix it automatically.\n\npisama.ai"),

        ("The Problem", "AI Agents Fail in Unpredictable Ways",
         """• Infinite loops - Agents repeat the same actions endlessly
• State corruption - Data becomes invalid mid-execution
• Persona drift - Agents forget their role and instructions
• Coordination deadlocks - Multi-agent systems get stuck
• Goal abandonment - Tasks silently fail to complete

Current debugging: Manual log reading for hours"""),

        ("The Solution", "PISAMA: Automated Agent Diagnostics",
         """Your Agent System (LangGraph | AutoGen | CrewAI | n8n)
                              ↓
                         PISAMA
                     26 Detectors
                     Auto-Healing
                              ↓
        Detection Alerts  |  Root Cause Analysis  |  Code Fixes"""),

        ("What We've Built", "Production-Ready Platform",
         """• Detection Engine - 26 failure detectors (LIVE)
• Self-Healing - AI-powered fix suggestions (LIVE)
• Dashboard - 11 views, real-time monitoring (LIVE)
• SDK - Python, 4 framework integrations (LIVE)
• CLI - CI/CD ready, MCP server (LIVE)
• API - REST, multi-tenant (LIVE)
• Demo - 4 interactive scenarios (LIVE)"""),

        ("Detection Engine", "26 Failure Mode Detectors",
         """BEHAVIORAL PATTERNS:
• Loop Detection - Structural, hash, and semantic analysis
• State Corruption - Domain validation, velocity anomalies
• Persona Drift - Role consistency scoring (5 role types)
• Coordination Failures - Circular delegation, ignored messages

SAFETY & SECURITY:
• Hallucination Detection - Grounding scores, citation verification
• Prompt Injection - 25+ attack patterns, 13 jailbreak signatures"""),

        ("Detection Engine (cont.)", "Resource, Performance & Task",
         """RESOURCE & PERFORMANCE:
• Context Overflow - Token warnings at 70%, 85%, 95%
• Cost Analysis - 25+ LLM models with 2025 pricing
• Latency Tracking - Per-span millisecond precision

TASK & WORKFLOW:
• Task Derailment - Off-topic deviation detection
• Specification Mismatch - Output vs requirements comparison
• Goal Abandonment - Incomplete task chain detection
• Quality Gate Bypass - Validation rule evasion"""),

        ("Detection Accuracy", "MAST 14-Mode Testing Framework",
         """System Design (F1-F5): Spec, decomposition, loops, tools, workflow → 92%
Inter-Agent (F6-F10): Derailment, context, withholding, coordination → 89%
Verification (F11-F14): Corruption, persona, quality gates, completion → 94%

DETECTION METHODS:
• Structural pattern matching
• Semantic similarity (pgvector embeddings)
• LLM-as-Judge with tier escalation"""),

        ("Self-Healing Engine", "AI-Powered Fix Suggestions",
         """EXAMPLE DETECTION:
Infinite loop between Agent1 and Agent2
Confidence: 94%

SUGGESTED FIX: Add circuit breaker
- Tracks consecutive failures
- Opens circuit after max_failures
- Auto-resets after timeout
- Prevents cascade failures"""),

        ("Fix Types", "Automated Remediation Options",
         """• Retry Limit - Add configurable retry counters
• Exponential Backoff - Progressive wait times
• Circuit Breaker - Stop after N failures
• State Validation - Add validation wrappers
• Schema Enforcement - Enforce data contracts
• Prompt Reinforcement - 3 levels: light, moderate, aggressive
• Role Boundary - Enforce agent responsibilities
• Timeout Addition - Prevent infinite waits
• Checkpoint Recovery - Resume from last good state"""),

        ("Dashboard", "11 Interactive Views",
         """• Dashboard - High-level metrics, recent detections
• Agent Orchestration - Network graph of agent communication
• Agent Details - Deep-dive into individual agents
• Traces - Execution history with filtering
• Detections - All failures by type/severity
• Agent Forensics - Paste-and-analyze debugging
• Testing - 14-mode accuracy metrics
• Chaos Engineering - Failure injection testing
• Replay - What-if scenario analysis
• Regression - Baseline comparison testing
• Demo - Interactive failure scenarios"""),

        ("Agent Forensics", "Root Cause Analysis in Seconds",
         """BEFORE PISAMA:
"The agent failed after 2 hours. I spent 4 hours reading logs."

WITH PISAMA:
1. Paste your trace (LangSmith, OpenTelemetry, or raw JSON)
2. Click "Diagnose"
3. See:
   • Primary failure with confidence score
   • Root cause chain with evidence
   • Affected spans highlighted
   • Suggested fix with code diff"""),

        ("Interactive Demo", "4 Failure Scenarios",
         """Healthy Workflow - Normal execution, all agents complete
Infinite Loop - Agents 1-2-3 repeat, detected at 94% confidence
State Corruption - Data degrades mid-execution
Coordination Deadlock - Agents waiting on each other

LIVE AT pisama.ai/demo
• Start/pause simulation
• Watch real-time detection alerts
• See agent activity feed"""),

        ("Python SDK", "3-Line Integration",
         """from mao_testing import MAOTracer

tracer = MAOTracer(api_key="your-key")

with tracer.session("my-workflow") as session:
    result = run_my_agents()
    session.snapshot("final_output", result)

FEATURES:
• OpenTelemetry-based tracing
• Automatic batch export
• Conditional sampling rules
• State snapshots at decision points"""),

        ("Framework Integrations", "Works With Your Stack",
         """• LangGraph - Full - decorator-based node tracing
• AutoGen - Full - agent + conversation tracing
• CrewAI - Full - crew, task, and agent tracing
• n8n - Full - async polling + webhooks
• LangChain - Planned
• OpenAI Assistants - Planned
• AWS Bedrock Agents - Planned"""),

        ("CLI & DevOps", "Production-Ready Tooling",
         """# Analyze a specific trace
mao debug trace-abc123

# Get fix suggestions
mao fix detection-xyz789 --apply

# Watch for new detections in real-time
mao watch --severity high

# CI/CD integration with exit codes
mao ci --threshold 0.95

FEATURES: Secure credentials, JSON output, CI/CD exit codes"""),

        ("MCP Server", "Claude Code Integration",
         """CONFIG (~/.claude/mcp.json):
{
  "mcpServers": {
    "pisama": { "command": "mao", "args": ["mcp", "serve"] }
  }
}

AVAILABLE TOOLS:
• mao_analyze_trace - Analyze specific traces
• mao_get_detections - Query detections
• mao_get_fix_suggestions - Get code fixes
• mao_get_trace - Full trace details

Rate limited, read-only, audit logged"""),

        ("API", "RESTful Multi-Tenant Architecture",
         """CORE ENDPOINTS:
POST /traces/ingest - Accept OTEL traces (async)
GET  /detections - List with filtering
POST /diagnose/why-failed - Root cause analysis
POST /diagnose/quick-check - Fast assessment
GET  /detections/{id}/fixes - Fix suggestions
POST /chaos/sessions - Chaos testing

ENTERPRISE: Clerk auth, per-tenant rate limiting, API keys, webhooks"""),

        ("Technical Architecture", "Production Infrastructure",
         """FRONTEND: Next.js 16 | React 18 | TailwindCSS

BACKEND: FastAPI | SQLAlchemy | PostgreSQL | pgvector

COMPONENTS:
• Detection Engine (26 algorithms)
• Ingestion Pipeline (4 formats)
• Self-Healing Engine (9 fix types)"""),

        ("Trace Ingestion", "Multi-Format Support",
         """• LangSmith - JSONL/JSON, run types, session mapping
• OpenTelemetry - OTEL + OTLP, hierarchical spans
• n8n - Workflow logs, node execution
• Raw JSON - Custom field mapping
• Universal - Framework-agnostic abstraction

PIPELINE FEATURES:
Async buffer, backpressure, token counting, PII sanitization"""),

        ("Competitive Landscape", "PISAMA vs Alternatives",
         """                    PISAMA  LangSmith  Arize  Braintrust
Multi-agent focus      YES      Limited    No       No
Self-healing fixes     YES      No         No       No
Loop detection         YES      No         No       No
Coordination analysis  YES      No         No       No
n8n integration        YES      No         No       No
MCP server             YES      No         No       No

OUR MOAT: Purpose-built for multi-agent systems"""),

        ("Roadmap", "What's Next",
         """Q1 2025:
• LangChain integration
• OpenAI Assistants integration
• Enhanced hallucination detection

Q2 2025:
• AWS Bedrock Agents integration
• Dify / Flowise integrations
• Auto-fix application (with approval)

Q3 2025:
• Trace replay & simulation
• Chaos engineering automation
• Team collaboration features"""),

        ("Pricing", "Simple, Usage-Based",
         """Free ($0) - 1K traces/mo, basic detections
Pro ($29/mo) - 50K traces, all detectors, fixes
Team ($99/mo) - 500K traces, team dashboard, API
Enterprise (Custom) - Unlimited, SLA, dedicated support

ALL TIERS INCLUDE:
• All 26 detection algorithms
• Dashboard access
• SDK and CLI"""),

        ("Getting Started", "5 Minutes to First Detection",
         """# 1. Install SDK
pip install mao-testing

# 2. Configure
mao config init

# 3. Add to your agent
from mao_testing import MAOTracer
tracer = MAOTracer()
with tracer.session("my-workflow") as s:
    result = my_agent.run()

# 4. View results
open https://pisama.ai/dashboard"""),

        ("Live Demo", "See PISAMA in Action",
         """pisama.ai/demo

• Select a failure scenario
• Watch agents execute in real-time
• See detection alerts fire
• Explore fix suggestions

OR TRY WITH YOUR OWN TRACES:
1. Go to pisama.ai/diagnose
2. Paste your LangSmith/OTEL trace
3. Get instant root cause analysis"""),

        ("Contact", None,
         """Website: pisama.ai
Demo: pisama.ai/demo
Docs: pisama.ai/docs

Ready to stop debugging agent failures manually?

Start your free trial today."""),
    ]

    for args in slides:
        add_slide(prs, *args)

    return prs

if __name__ == "__main__":
    prs = create_presentation()
    prs.save("PISAMA_Product_Editable.pptx")
    print("Created: PISAMA_Product_Editable.pptx")
